# Build the operator-grouped PDF-report status spreadsheet AND rebuild the two deck
# table slides (0715 Data Subcommittee Report). Single data source: GROUPS below.
# xlsx keeps full names + separate Vehicle/Reg columns; the deck merges Vehicle+Reg into
# one column and uses compact display names so every cell stays single-line at 18 pt.
from pathlib import Path

from openpyxl import Workbook
from openpyxl.styles import Alignment, Border, Font, PatternFill, Side
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
XLSX = ROOT / "pdf_report_workspace" / "pdf_report_status.xlsx"
PPTX = ROOT / "monthly_presentation" / "20260715" / "20260715Data Subcommittee Report v1.0.pptx"

NEEDED = "Needed (no data yet)"
REQUESTED = "Requested (no data yet)"  # operating data asked of the operator (with the sent PDF)
# Round-1 "PDF sent" markers (user-reported 2026-07-10): report delivered to the operator.
SENT_R1 = {("SJG (DP World II)", "CMZ6260"), ("Port Express (DP World II)", "YN75NMA"),
           ("Nestlé", "EV73SAL"), ("Nestlé", "YK73WFN"),
           ("William Jackson Food", "EX74JXW"), ("William Jackson Food", "YN25RSY")}
# WU70GLV's 2025-09→11 legs carry a WJF round-robin token in SRF, but the vehicle is DP
# World's and was never lent to WJF (user-confirmed SRF data error, 2026-07-10) — it is
# therefore ONE continuous DP World trial here, and only YT21EFD is WJF's comparator.
# (operator display name, [(vehicle_full, vehicle_short, reg, trial_type, diesel_comparator, period, pdf_ok)])
# pdf_ok: True = PDF generated, False = not generated, None = PLANNED trial (data collection
# upcoming, reg often TBD) — planned rows appear in the xlsx only, never on the deck slides.
# DP World family first (I = direct 2025 trials; II = the 2026 round via its subsidiaries
# SJG / Port Express), then alphabetical.
GROUPS = [
    ("DP World I", [
        ("Scania P-series BEV", "Scania BEV", "EX74JXW", "Round-robin", "WU70GLV", "2025-07-07 – 2025-08-22", True),
        ("DAF XF 450 (diesel)", "DAF XF 450", "WU70GLV", "Round-robin", "—", "2025-06-26 – 2025-12-11", False),
        ("MAN BEV (new model)", "MAN BEV", "TBD", "TBD", "WU70GLV", "Not started", None),
        ("Volvo BEV (telematics fault)", "Volvo BEV", "TBD", "BYO", "WU70GLV", "Not started", None),
        ("Diesel comparators ×7 (2+2+3)", "Diesel ×7", "TBD", "BYO", "—", "Not started", None),
    ]),
    ("Port Express (DP World II)", [
        ("Mercedes-Benz eActros 600", "eActros 600", "YN75NMA", "Round-robin", REQUESTED, "2026-01-29 – 2026-04-08", True),
    ]),
    ("SJG (DP World II)", [
        ("Volvo FH Electric", "Volvo FH", "CMZ6260", "Round-robin", REQUESTED, "2026-02-06 – 2026-04-07", True),
    ]),
    ("Co-Op", [
        ("Scania BEV (health check TBD)", "Scania BEV", "MK15BEV", "BYO", "2 planned (TBD)", "Not started", None),
        ("Diesel comparator (first)", "Diesel", "TBD", "BYO", "—", "Not started", None),
        ("Diesel comparator (second)", "Diesel", "TBD", "BYO", "—", "Not started", None),
    ]),
    ("HTL (Howard Tenens)", [
        ("Volvo FH Electric", "Volvo FH", "CMZ6260", "Round-robin", NEEDED, "2026-04-27 – 2026-07-01", True),
        ("Mercedes-Benz eActros 600", "eActros 600", "YN75NMA", "Round-robin", NEEDED, "2026-04-17 – 2026-06-26", True),
        ("Diesel comparator", "Diesel", "TBD", "BYO", "—", "Not started", None),
    ]),
    ("John Lewis Partnership", [
        ("Volvo FM Electric", "Volvo FM", "KY24LHT", "BYO", NEEDED, "2024-06-21 – 2025-03-20", True),
        ("Scania P-series BEV", "Scania BEV", "EX74JXY", "BYO", NEEDED, "2025-04-11 – 2026-05-03", True),
        ("DAF XD Electric", "DAF XD", "LN25NKE", "Round-robin", NEEDED, "2025-09-01 – 2025-12-29", True),
        ("Volvo FH Electric", "Volvo FH", "CMZ6260", "Round-robin", NEEDED, "2025-10-30 – 2025-12-23", True),
    ]),
    ("Knowles", [
        ("Volvo FM Electric", "Volvo FM", "AV24LXJ", "BYO", NEEDED, "2024-06-11 – 2026-07-03", True),
        ("Volvo FM Electric", "Volvo FM", "AV24LXK", "BYO", NEEDED, "2024-06-11 – 2026-07-03", True),
        ("Volvo FM Electric", "Volvo FM", "AV24LXL", "BYO", NEEDED, "2024-06-11 – 2026-07-03", True),
        ("Diesel comparator", "Diesel", "TBD", "BYO", "—", "Not started", None),
    ]),
    ("Nestlé", [
        ("Volvo FM Electric", "Volvo FM", "EV73SAL", "BYO", REQUESTED, "2024-06-12 – 2026-07-05", True),
        ("Volvo FM Electric", "Volvo FM", "YK73WFN", "BYO", REQUESTED, "2024-06-12 – 2026-07-03", True),
    ]),
    ("Welch's", [
        ("Renault E-Tech D Wide", "Renault D Wide", "T88RNW", "BYO", NEEDED, "2024-06-11 – 2026-07-03", True),
        ("Renault Trucks D Wide Z.E.", "Renault D Wide ZE", "N88GNW", "BYO", NEEDED, "2024-10-11 – 2026-07-03", True),
        ("Renault E-Tech T", "Renault E-Tech T", "TA70WTL", "BYO", NEEDED, "2025-05-03 – 2026-07-03", True),
        ("DAF XD Electric", "DAF XD", "LN25NKE", "Round-robin", NEEDED, "2026-01-14 – 2026-04-07", True),
        ("Scania P-series BEV", "Scania BEV", "EX74JXW", "Round-robin", NEEDED, "2026-02-26 – 2026-04-29", True),
    ]),
    ("William Jackson Food", [
        ("Scania P-series BEV", "Scania BEV", "EX74JXW", "Round-robin", "YT21EFD", "2025-10-11 – 2025-11-18", True),
        ("Mercedes-Benz eActros 600", "eActros 600", "YN25RSY", "Round-robin", "YT21EFD", "2025-10-21 – 2025-11-18", True),
        ("Scania P410 (diesel)", "Scania P410", "YT21EFD", "BYO", "—", "2025-08-30 – 2026-07-04", False),
    ]),
    ("WS", [
        ("DAF XD Electric", "DAF XD", "LN25NKE", "Round-robin", NEEDED, "2026-04-16 – 2026-07-03", True),
    ]),
]
# xlsx: 5 fixed columns + REPEATABLE per-round groups (Trial period / PDF generated /
# PDF sent) so each reporting round is recorded separately; Round 2 is pre-created blank.
XL_FIXED = ["Operator", "Vehicle", "Reg", "Trial type", "Diesel comparator"]
XL_ROUND_SUB = ["Trial period", "PDF generated", "PDF sent"]
XL_ROUNDS = ["Round 1 (2026-07-10)", "Round 2"]
PPT_HEADERS = ["Operator", "Vehicle", "Trial type", "Diesel comparator", "Trial period", "PDF report"]

NAVY = RGBColor(0x1F, 0x49, 0x7D)
ACCENTS = [RGBColor(0x4F, 0x81, 0xBD), RGBColor(0xC0, 0x50, 0x4D), RGBColor(0x9B, 0xBB, 0x59),
           RGBColor(0x80, 0x64, 0xA2), RGBColor(0x4B, 0xAC, 0xC6), RGBColor(0xF7, 0x96, 0x46), NAVY]
PPT_ACCENTS = [ACCENTS[0], ACCENTS[1], ACCENTS[3], ACCENTS[4], ACCENTS[5], ACCENTS[6]]
GREEN_TXT = RGBColor(0x4F, 0x62, 0x28)
RED_TXT = RGBColor(0x94, 0x37, 0x34)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x26, 0x26, 0x26)


def tint(c, f):
    return RGBColor(*(int(v + (255 - v) * f) for v in (c[0], c[1], c[2])))


def hexs(c):
    return f"{c[0]:02X}{c[1]:02X}{c[2]:02X}"


# ── xlsx: 5 fixed columns + repeatable round groups (two-level header) ──
wb = Workbook()
ws = wb.active
ws.title = "PDF report status"
ROUND_COLS = len(XL_ROUND_SUB)
widths = [33, 36, 13, 17, 27] + [30, 18, 18] * len(XL_ROUNDS)
for j, w in enumerate(widths, start=1):
    ws.column_dimensions[ws.cell(row=1, column=j).column_letter].width = w
thin = Side(style="thin", color="BFBFBF")
border = Border(left=thin, right=thin, top=thin, bottom=thin)
# per-column accents: fixed 5, then per-round (period orange, generated navy, sent green)
ROUND_ACCENTS = [ACCENTS[5], ACCENTS[6], ACCENTS[2]]
col_accent = ACCENTS[:5] + ROUND_ACCENTS * len(XL_ROUNDS)


def _hdr(cell, text, fill):
    cell.value = text
    cell.font = Font(name="Arial", bold=True, color="FFFFFF", size=14)
    cell.fill = PatternFill("solid", fgColor=hexs(fill))
    cell.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)
    cell.border = border


for j, h in enumerate(XL_FIXED, start=1):
    _hdr(ws.cell(row=1, column=j), h, ACCENTS[j - 1])
    ws.merge_cells(start_row=1, start_column=j, end_row=2, end_column=j)
for ri, rname in enumerate(XL_ROUNDS):
    c0 = len(XL_FIXED) + 1 + ri * ROUND_COLS
    _hdr(ws.cell(row=1, column=c0), rname, NAVY if ri == 0 else tint(NAVY, 0.35))
    ws.merge_cells(start_row=1, start_column=c0, end_row=1, end_column=c0 + ROUND_COLS - 1)
    for k, sub in enumerate(XL_ROUND_SUB):
        _hdr(ws.cell(row=2, column=c0 + k), sub, ROUND_ACCENTS[k])
ws.freeze_panes = "A3"
r = 3
for op, rows in GROUPS:
    g0 = r
    for veh, _vs, reg, ttype, comp, period, ok in rows:
        light = 0.80 if (r % 2 == 0) else 0.90
        # round 1 filled from the current batch; PDF sent left for manual upkeep; round 2 blank.
        # PDF-generated states are WORDS (no symbols, user 2026-07-10): Done / Not needed
        # (diesel comparators - no partner PDF planned) / No data yet (planned rows) /
        # "To be done" (pass a string in GROUPS when data exists but the PDF is pending).
        gen = {True: "Done", False: "Not needed", None: "No data yet"}.get(ok, str(ok))
        sent = "Sent" if (op, reg) in SENT_R1 else ""
        vals = [op, veh, reg, ttype, comp, period, gen, sent] + [""] * ROUND_COLS
        for j, v in enumerate(vals, start=1):
            c = ws.cell(row=r, column=j, value=v if v != "" else None)
            bold = (j == 3 and v != "TBD") or (j == 5 and v not in (NEEDED, REQUESTED, "—", "Not needed") and "planned" not in v)
            if ok is None:
                colour = "7F7F7F"
            else:
                colour = "262626"
            if j == 7:
                colour = {"Done": "4F6228", "To be done": "BF8F00"}.get(v, "595959")
            if j == 8 and v == "Sent":
                colour = "4F6228"
            c.font = Font(name="Arial", size=14, bold=bold, color=colour)
            c.fill = PatternFill("solid", fgColor=hexs(tint(col_accent[j - 1], light)))
            c.alignment = Alignment(horizontal="left" if j in (1, 2, 6, 9) else "center",
                                    vertical="center")
            c.border = border
        r += 1
    if len(rows) > 1:
        ws.merge_cells(start_row=g0, start_column=1, end_row=r - 1, end_column=1)
        ws.cell(row=g0, column=1).alignment = Alignment(horizontal="left", vertical="center")
NOTES = [
    "PDF generated / PDF sent — Done: report generated; Sent: report delivered; To be done: "
    "data available, report pending; No data yet: data collection not started; Not needed: no "
    "partner PDF planned (diesel comparator vehicles).",
    "Diesel comparator — a reg (e.g. YT21EFD): comparator data exists; Requested (no data yet): "
    "operating data requested from the operator; Needed (no data yet): planned, not yet requested; "
    "—: the row is itself a comparator.",
]
for k, note in enumerate(NOTES):
    c = ws.cell(row=r + 1 + k, column=1, value=note)
    c.font = Font(name="Arial", size=14, italic=True)
wb.save(XLSX)
print("xlsx saved:", XLSX)

# ───────────────────────── pptx (6 columns, Vehicle·Reg merged) ─────────────────────────
SLIDE_GROUPS = [GROUPS[:6], GROUPS[6:]]  # 13 + 11 data rows
TITLES = ["Partner PDF Reports — Coverage by Operator (1/2)",
          "Partner PDF Reports — Coverage by Operator (2/2)"]
COL_W = [1.75, 3.75, 1.60, 1.70, 2.80, 1.10]  # sums to 12.70
FONT, SZ = "Arial", 18

prs = Presentation(PPTX)


def set_cell(cell, runs, colour, fill, center=False, wrap=False):
    """runs: [(text, bold), ...] rendered in one paragraph."""
    cell.fill.solid()
    cell.fill.fore_color.rgb = fill
    cell.margin_left, cell.margin_right = Inches(0.05), Inches(0.04)
    cell.margin_top = cell.margin_bottom = Inches(0.01)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.text_frame.word_wrap = wrap
    p = cell.text_frame.paragraphs[0]
    for text, bold in runs:
        run = p.add_run()
        run.text = text
        f = run.font
        f.name, f.size, f.bold = FONT, Pt(SZ), bold
        f.color.rgb = colour
    if center:
        p.alignment = PP_ALIGN.CENTER


for slide, groups, title in zip([prs.slides[1], prs.slides[2]], SLIDE_GROUPS, TITLES):
    # deck shows ACTIVE trials only — drop planned rows (ok=None) and then-empty groups
    groups = [(op, [row for row in oprows if row[6] is not None]) for op, oprows in groups]
    groups = [(op, oprows) for op, oprows in groups if oprows]
    for sh in list(slide.shapes):
        if sh.has_table:
            sh._element.getparent().remove(sh._element)
    slide.shapes.title.text = title
    trun = slide.shapes.title.text_frame.paragraphs[0].runs[0]
    trun.font.color.rgb = NAVY
    trun.font.name = FONT

    rows = [(op, row) for op, oprows in groups for row in oprows]
    n_rows = len(rows) + 1
    left = Inches((13.333 - sum(COL_W)) / 2)
    # single-row groups with a long operator name wrap to 2 lines -> taller row
    tall = set()
    r0 = 1
    for op, oprows in groups:
        if len(oprows) == 1 and len(op) > 14:
            tall.add(r0)
        r0 += len(oprows)
    height = 0.62 + 0.34 * len(rows) + 0.28 * len(tall)
    tbl = slide.shapes.add_table(n_rows, 6, left, Inches(1.10), Inches(sum(COL_W)), Inches(height)).table
    tbl.first_row = False
    tbl.horz_banding = False
    for j, w in enumerate(COL_W):
        tbl.columns[j].width = Inches(w)
    tbl.rows[0].height = Inches(0.62)
    for i in range(1, n_rows):
        tbl.rows[i].height = Inches(0.62 if i in tall else 0.34)
    for j, h in enumerate(PPT_HEADERS):
        set_cell(tbl.cell(0, j), [(h, True)], WHITE, PPT_ACCENTS[j], center=True, wrap=True)

    op_spans, r0 = [], 1
    for op, oprows in groups:
        if len(oprows) > 1:
            op_spans.append((r0, r0 + len(oprows) - 1))
        r0 += len(oprows)
    op_merged = {(i, 0) for lo, hi in op_spans for i in range(lo + 1, hi + 1)}
    op_fill = {i: tint(PPT_ACCENTS[0], 0.85) for lo, hi in op_spans for i in range(lo, hi + 1)}

    for i, (op, (_vf, veh_s, reg, ttype, comp, period, ok)) in enumerate(rows, start=1):
        light = 0.80 if i % 2 == 0 else 0.90
        cells = [
            ([(op if (i, 0) not in op_merged else "", False)],
             BLACK, op_fill.get(i, tint(PPT_ACCENTS[0], light)), False, True),
            ([(veh_s + " · ", False), (reg, True)], BLACK, tint(PPT_ACCENTS[1], light), False, False),
            ([(ttype, False)], BLACK, tint(PPT_ACCENTS[2], light), True, False),
            ([({NEEDED: "Awaiting data", REQUESTED: "Requested"}.get(comp, comp), comp not in (NEEDED, REQUESTED, "—"))],
             BLACK, tint(PPT_ACCENTS[3], light), True, False),
            ([(period.replace(" – ", "–"), False)], BLACK, tint(PPT_ACCENTS[4], light), False, False),
            ([("✓" if ok else "✗", True)], GREEN_TXT if ok else RED_TXT, tint(PPT_ACCENTS[5], light), True, False),
        ]
        for j, (runs, colour, fill, center, wrap) in enumerate(cells):
            set_cell(tbl.cell(i, j), runs, colour, fill, center=center, wrap=wrap)
    for lo, hi in op_spans:
        tbl.cell(lo, 0).merge(tbl.cell(hi, 0))

prs.save(PPTX)
print("pptx saved:", PPTX)
